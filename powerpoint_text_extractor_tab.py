import os 
import pptx
import nltk
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from langdetect import detect, DetectorFactory
from langdetect.lang_detect_exception import LangDetectException

# Ensure consistent language detection
DetectorFactory.seed = 0

# Lazy initialization of NLTK data
def ensure_nltk_data():
    """
    Check if the 'punkt' tokenizer is available. Download it if it's missing.
    """
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt')

class PowerPointTextExtractorTab:
    def __init__(self, parent):
        self.frame = ttk.Frame(parent, padding=10)
        self.build_widgets()   

    def build_widgets(self):
        # PowerPoint file selection
        ttk.Label(self.frame, text="Select PowerPoint File:").grid(row=0, column=0, padx=5, pady=5, sticky="w")
        self.file_entry = ttk.Entry(self.frame, width=50)
        self.file_entry.grid(row=0, column=1, padx=5, pady=5)
        ttk.Button(self.frame, text="Browse", command=self.browse_file).grid(row=0, column=2, padx=5, pady=5)

        # Output folder selection
        ttk.Label(self.frame, text="Select Output Folder:").grid(row=1, column=0, padx=5, pady=5, sticky="w")
        self.output_entry = ttk.Entry(self.frame, width=50)
        self.output_entry.grid(row=1, column=1, padx=5, pady=5)
        ttk.Button(self.frame, text="Browse", command=self.browse_output_folder).grid(row=1, column=2, padx=5, pady=5)

        # Language selection
        ttk.Label(self.frame, text="Select Language:").grid(row=2, column=0, padx=5, pady=5, sticky="w")
        self.lang_var = tk.StringVar(value="de")   # Default language is German
        lang_menu = ttk.OptionMenu(self.frame, self.lang_var, "de", "de", "en")
        lang_menu.grid(row=2, column=1, padx=5, pady=5, sticky="w")

        # Checkbox for separating sentences and words
        self.separate_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(self.frame, text="Separate Sentences and Words", variable=self.separate_var).grid(row=3, column=1, padx=5, pady=5, sticky="w")

        # Checkbox for filtering numeric-only lines
        self.filter_numeric_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(self.frame, text="Filter Lines with Only Numbers", variable=self.filter_numeric_var).grid(row=4, column=1, padx=5, pady=5, sticky="w")
        
        # Minimum characters for detection entry
        ttk.Label(self.frame, text="Minimum Characters for Detection:").grid(row=5, column=0, padx=5, pady=5, sticky="w")
        self.threshold_var = tk.StringVar(value="5")
        ttk.Entry(self.frame, textvariable=self.threshold_var, width=10).grid(row=5, column=1, padx=5, pady=5, sticky="w")
        
        # Checkbox to enable or disable language detection
        self.detect_lang_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(self.frame, text="Enable Language Detection", variable=self.detect_lang_var).grid(row=6, column=1, padx=5, pady=5, sticky="w")
        
        # Checkbox to show extraction statistics (word and sentence counts)
        self.show_counts_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(self.frame, text="Show Extraction Statistics", variable=self.show_counts_var).grid(row=7, column=1, padx=5, pady=5, sticky="w")
        
        # Extract button
        ttk.Button(self.frame, text="Extract", command=self.start_extraction).grid(row=8, column=1, padx=5, pady=10)

    def browse_file(self):
        filename = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if filename:
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, filename)

    def browse_output_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, folder)

    # Extracts text from a PowerPoint file
    def extract_text_from_pptx(self, pptx_path):
        presentation = pptx.Presentation(pptx_path)
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text.strip())
        return texts

    def filter_text_by_language(self, texts, selected_lang, filter_numeric=False, threshold=10, perform_detection=True):
        filtered_sentences = []
        filtered_words = []
        for text in texts:
            lines = text.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                if filter_numeric and line.isdigit():
                    continue
                if perform_detection and len(line) >= threshold:
                    try:
                        if detect(line) != selected_lang:
                            continue
                    except LangDetectException:
                        continue
                if line and line[-1] in ".!?":
                    filtered_sentences.append(line)
                else:
                    filtered_words.append(line)
        return filtered_sentences, filtered_words

    # Save the extracted text to a file
    def save_to_file(self, filename, data, add_spacing=False):
        with open(filename, "w", encoding="utf-8") as file:
            for item in data:
                file.write(item + "\n")
                if add_spacing:
                    file.write("\n")

    def process_file(self, pptx_path, output_folder, selected_lang, separate_sentences_words, filter_numeric, threshold, language_detection, show_counts):
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
        texts = self.extract_text_from_pptx(pptx_path)
        sentences, words = self.filter_text_by_language(
            texts, 
            selected_lang, 
            filter_numeric=filter_numeric, 
            threshold=threshold, 
            perform_detection=language_detection
        )
        if separate_sentences_words:
            self.save_to_file(os.path.join(output_folder, "sentences.txt"), sentences, add_spacing=True)
            self.save_to_file(os.path.join(output_folder, "words.txt"), words)
        else:
            self.save_to_file(os.path.join(output_folder, "text.txt"), sentences + words)
        
        # Prepare the final message with optional extraction statistics.
        message = "Extraction completed. Check the output folder for results."
        if show_counts:
            message += f"\n\nSentences detected: {len(sentences)}\nWords detected: {len(words)}"
        messagebox.showinfo("Success", message)

    def start_extraction(self):
        ensure_nltk_data()
        
        pptx_path = self.file_entry.get()
        output_folder = self.output_entry.get()
        selected_lang = self.lang_var.get()
        separate = self.separate_var.get()
        filter_numeric = self.filter_numeric_var.get()
        try:
            threshold = int(self.threshold_var.get())
        except ValueError:
            messagebox.showerror("Error", "Minimum characters threshold must be an integer.")
            return
        
        language_detection = self.detect_lang_var.get()
        show_counts = self.show_counts_var.get()
        
        if not pptx_path:
            messagebox.showerror("Error", "Please select a PowerPoint file.")
            return
        if not output_folder:
            messagebox.showerror("Error", "Please select an output folder.")
            return
        self.process_file(pptx_path, output_folder, selected_lang, separate, filter_numeric, threshold, language_detection, show_counts)
